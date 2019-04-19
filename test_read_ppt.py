from pptx import Presentation

prs = Presentation('test.pptx')

for slide in prs.slides:
    print('slide', slide.name)
    print("(%s, %s)" % (slide.width, slide.height))
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        print(shape.text)
        print("(%s, %s)" % (shape.width, shape.height))

        # for paragraph in shape.text_frame.paragraphs:
        #     for run in paragraph.runs:
        #         print(run.text, end='_')

